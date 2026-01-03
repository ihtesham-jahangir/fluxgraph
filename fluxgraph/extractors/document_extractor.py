"""
FluxGraph Document Extractor

Extract text and data from various document formats:
- PDF files
- Word documents (DOCX)
- Excel spreadsheets (XLSX)
- PowerPoint presentations (PPTX)
- Text files
- Markdown files
- CSV files
"""

import logging
from typing import Dict, List, Optional, Union, Any
from pathlib import Path
import mimetypes

try:
    import httpx
except ImportError:
    httpx = None


logger = logging.getLogger(__name__)


class DocumentContent:
    """Container for extracted document content"""

    def __init__(
        self,
        text: str,
        metadata: Optional[Dict] = None,
        pages: Optional[List[str]] = None,
        tables: Optional[List] = None,
        images: Optional[List] = None
    ):
        self.text = text
        self.metadata = metadata or {}
        self.pages = pages or []
        self.tables = tables or []
        self.images = images or []

    def to_dict(self) -> Dict:
        return {
            "text": self.text,
            "metadata": self.metadata,
            "page_count": len(self.pages),
            "table_count": len(self.tables),
            "image_count": len(self.images)
        }


class DocumentExtractor:
    """
    Universal document extractor

    Example:
        extractor = DocumentExtractor()

        # From file path
        content = extractor.extract_from_file("document.pdf")
        print(content.text)

        # From URL
        content = await extractor.extract_from_url("https://example.com/doc.pdf")

        # Auto-detect format
        content = extractor.extract("path/to/document.pdf")
    """

    def __init__(self):
        self.extractors = {
            '.pdf': self._extract_pdf,
            '.docx': self._extract_docx,
            '.doc': self._extract_doc,
            '.xlsx': self._extract_xlsx,
            '.xls': self._extract_xls,
            '.pptx': self._extract_pptx,
            '.txt': self._extract_txt,
            '.md': self._extract_txt,
            '.csv': self._extract_csv,
            '.json': self._extract_json,
            '.xml': self._extract_xml,
            '.html': self._extract_html
        }

    def extract_from_file(self, file_path: Union[str, Path]) -> DocumentContent:
        """
        Extract content from file

        Args:
            file_path: Path to file

        Returns:
            DocumentContent object
        """
        file_path = Path(file_path)

        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        extension = file_path.suffix.lower()
        extractor = self.extractors.get(extension)

        if extractor is None:
            raise ValueError(f"Unsupported file format: {extension}")

        logger.info(f"📄 Extracting {extension} file: {file_path}")
        return extractor(file_path)

    async def extract_from_url(self, url: str) -> DocumentContent:
        """
        Download and extract content from URL

        Args:
            url: URL to document

        Returns:
            DocumentContent object
        """
        if httpx is None:
            raise ImportError("httpx required for URL extraction. Install with: pip install httpx")

        async with httpx.AsyncClient() as client:
            response = await client.get(url, follow_redirects=True)
            response.raise_for_status()

            # Detect file type from Content-Type or URL
            content_type = response.headers.get("content-type", "")
            extension = self._detect_extension(url, content_type)

            # Save to temp file
            import tempfile
            with tempfile.NamedTemporaryFile(suffix=extension, delete=False) as tmp:
                tmp.write(response.content)
                tmp_path = Path(tmp.name)

            try:
                return self.extract_from_file(tmp_path)
            finally:
                tmp_path.unlink()  # Clean up

    def extract(self, source: Union[str, Path]) -> DocumentContent:
        """
        Auto-detect and extract

        Args:
            source: File path or content string

        Returns:
            DocumentContent object
        """
        # Check if it's a file path
        if isinstance(source, (str, Path)):
            path = Path(source)
            if path.exists():
                return self.extract_from_file(path)

        raise ValueError("Invalid source. Provide a file path.")

    def _extract_pdf(self, file_path: Path) -> DocumentContent:
        """Extract text from PDF"""
        try:
            from pypdf import PdfReader
        except ImportError:
            raise ImportError("pypdf required for PDF extraction. Install with: pip install pypdf")

        reader = PdfReader(file_path)

        text_parts = []
        pages = []

        for i, page in enumerate(reader.pages):
            page_text = page.extract_text()
            text_parts.append(page_text)
            pages.append(page_text)

        metadata = {
            "page_count": len(reader.pages),
            "file_name": file_path.name,
            "file_size": file_path.stat().st_size
        }

        # Add PDF metadata if available
        if reader.metadata:
            metadata.update({
                "title": reader.metadata.get("/Title", ""),
                "author": reader.metadata.get("/Author", ""),
                "subject": reader.metadata.get("/Subject", ""),
                "creator": reader.metadata.get("/Creator", "")
            })

        return DocumentContent(
            text="\n\n".join(text_parts),
            metadata=metadata,
            pages=pages
        )

    def _extract_docx(self, file_path: Path) -> DocumentContent:
        """Extract text from DOCX"""
        try:
            from docx import Document
        except ImportError:
            raise ImportError("python-docx required. Install with: pip install python-docx")

        doc = Document(file_path)

        text_parts = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)

        # Extract tables
        tables = []
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            tables.append(table_data)

        metadata = {
            "paragraph_count": len(doc.paragraphs),
            "table_count": len(doc.tables),
            "file_name": file_path.name,
            "file_size": file_path.stat().st_size
        }

        # Core properties
        if hasattr(doc.core_properties, 'title'):
            metadata["title"] = doc.core_properties.title or ""
            metadata["author"] = doc.core_properties.author or ""
            metadata["subject"] = doc.core_properties.subject or ""

        return DocumentContent(
            text="\n\n".join(text_parts),
            metadata=metadata,
            tables=tables
        )

    def _extract_doc(self, file_path: Path) -> DocumentContent:
        """Extract text from DOC (legacy Word format)"""
        try:
            import textract
        except ImportError:
            raise ImportError("textract required for DOC extraction. Install with: pip install textract")

        text = textract.process(str(file_path)).decode('utf-8')

        return DocumentContent(
            text=text,
            metadata={"file_name": file_path.name}
        )

    def _extract_xlsx(self, file_path: Path) -> DocumentContent:
        """Extract data from Excel"""
        try:
            import openpyxl
        except ImportError:
            raise ImportError("openpyxl required. Install with: pip install openpyxl")

        workbook = openpyxl.load_workbook(file_path, data_only=True)

        text_parts = []
        tables = []

        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]

            text_parts.append(f"=== Sheet: {sheet_name} ===")

            sheet_data = []
            for row in sheet.iter_rows(values_only=True):
                row_data = [str(cell) if cell is not None else "" for cell in row]
                sheet_data.append(row_data)
                text_parts.append("\t".join(row_data))

            tables.append({
                "sheet": sheet_name,
                "data": sheet_data
            })

        metadata = {
            "sheet_count": len(workbook.sheetnames),
            "sheets": workbook.sheetnames,
            "file_name": file_path.name
        }

        return DocumentContent(
            text="\n".join(text_parts),
            metadata=metadata,
            tables=tables
        )

    def _extract_xls(self, file_path: Path) -> DocumentContent:
        """Extract data from Excel (legacy)"""
        try:
            import xlrd
        except ImportError:
            raise ImportError("xlrd required for XLS. Install with: pip install xlrd")

        workbook = xlrd.open_workbook(file_path)

        text_parts = []
        tables = []

        for sheet in workbook.sheets():
            text_parts.append(f"=== Sheet: {sheet.name} ===")

            sheet_data = []
            for row_idx in range(sheet.nrows):
                row_data = [str(cell.value) for cell in sheet.row(row_idx)]
                sheet_data.append(row_data)
                text_parts.append("\t".join(row_data))

            tables.append({
                "sheet": sheet.name,
                "data": sheet_data
            })

        return DocumentContent(
            text="\n".join(text_parts),
            metadata={"sheet_count": workbook.nsheets},
            tables=tables
        )

    def _extract_pptx(self, file_path: Path) -> DocumentContent:
        """Extract text from PowerPoint"""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx required. Install with: pip install python-pptx")

        prs = Presentation(file_path)

        text_parts = []
        pages = []

        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)

            slide_content = "\n".join(slide_text)
            text_parts.append(f"=== Slide {i + 1} ===\n{slide_content}")
            pages.append(slide_content)

        return DocumentContent(
            text="\n\n".join(text_parts),
            metadata={"slide_count": len(prs.slides), "file_name": file_path.name},
            pages=pages
        )

    def _extract_txt(self, file_path: Path) -> DocumentContent:
        """Extract text from TXT/MD files"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()

        return DocumentContent(
            text=text,
            metadata={"file_name": file_path.name, "file_size": file_path.stat().st_size}
        )

    def _extract_csv(self, file_path: Path) -> DocumentContent:
        """Extract data from CSV"""
        import csv

        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            reader = csv.reader(f)
            rows = list(reader)

        text = "\n".join([",".join(row) for row in rows])

        return DocumentContent(
            text=text,
            metadata={"row_count": len(rows), "file_name": file_path.name},
            tables=[{"data": rows}]
        )

    def _extract_json(self, file_path: Path) -> DocumentContent:
        """Extract data from JSON"""
        import json

        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)

        text = json.dumps(data, indent=2)

        return DocumentContent(
            text=text,
            metadata={"file_name": file_path.name}
        )

    def _extract_xml(self, file_path: Path) -> DocumentContent:
        """Extract text from XML"""
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            raise ImportError("beautifulsoup4 required. Install with: pip install beautifulsoup4")

        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()

        soup = BeautifulSoup(content, 'xml') or BeautifulSoup(content, 'html.parser')
        text = soup.get_text()

        return DocumentContent(
            text=text,
            metadata={"file_name": file_path.name}
        )

    def _extract_html(self, file_path: Path) -> DocumentContent:
        """Extract text from HTML"""
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            raise ImportError("beautifulsoup4 required. Install with: pip install beautifulsoup4")

        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()

        soup = BeautifulSoup(content, 'html.parser')

        # Remove script and style
        for element in soup(["script", "style"]):
            element.decompose()

        text = soup.get_text(separator="\n", strip=True)

        return DocumentContent(
            text=text,
            metadata={"file_name": file_path.name}
        )

    def _detect_extension(self, url: str, content_type: str) -> str:
        """Detect file extension from URL or content type"""
        # Try URL first
        from urllib.parse import urlparse
        path = urlparse(url).path
        if '.' in path:
            return Path(path).suffix

        # Try content type
        extension = mimetypes.guess_extension(content_type.split(';')[0])
        return extension or '.bin'


# Convenience functions

def extract_text(file_path: Union[str, Path]) -> str:
    """
    Quick text extraction

    Example:
        text = extract_text("document.pdf")
        print(text)
    """
    extractor = DocumentExtractor()
    content = extractor.extract_from_file(file_path)
    return content.text


async def extract_from_url(url: str) -> str:
    """
    Quick extraction from URL

    Example:
        text = await extract_from_url("https://example.com/doc.pdf")
    """
    extractor = DocumentExtractor()
    content = await extractor.extract_from_url(url)
    return content.text
